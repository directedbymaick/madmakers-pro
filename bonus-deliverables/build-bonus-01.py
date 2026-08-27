#!/usr/bin/env python3
"""
Build script for Bonus #1:
"Fiche Google Parfaite - 12 points pour la page 1 locale"
Carnet Plein(R) by Mad Makers

Generates 01-fiche-google-parfaite.pptx (16:9 widescreen).
Upload to Google Drive then right-click > "Ouvrir avec Google Slides"
- the file becomes a native, fully editable Google Slides document.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ===========================================================
# CARNET PLEIN(R) BRAND CHARTER
# ===========================================================
ENCRE        = RGBColor(0x0a, 0x0a, 0x0a)
ENCRE_SOFT   = RGBColor(0x1a, 0x1c, 0x18)
PAPIER       = RGBColor(0xf3, 0xf2, 0xee)
PAPIER_CHAUD = RGBColor(0xeb, 0xe7, 0xdc)
ACCENT       = RGBColor(0xe0, 0x54, 0x1b)
ACCENT_SOFT  = RGBColor(0xfa, 0xe2, 0xd5)  # lighter accent for backgrounds
INK          = RGBColor(0x1a, 0x1c, 0x18)
INK_DIM      = RGBColor(0x5a, 0x5d, 0x56)
PAPER_DIM    = RGBColor(0xa9, 0xa6, 0x9f)
RULE_LIGHT   = RGBColor(0xd5, 0xd2, 0xc9)
WHITE        = RGBColor(0xff, 0xff, 0xff)

FONT_BODY    = "Inter"
FONT_DISPLAY = "Inter"
FONT_MONO    = "Roboto Mono"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets")

# ===========================================================
# INIT
# ===========================================================
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ===========================================================
# LOW-LEVEL HELPERS
# ===========================================================
def blank_slide(bg_color=PAPIER):
    """New slide with solid background."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    return s


def text(slide, x, y, w, h, content, *, size=18, bold=False, color=ENCRE,
         font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.2, tracking=0):
    """Single-run text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing

    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def rich(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.3):
    """Text box with multiple runs (e.g. mixed bold / colors).
    runs = list of dicts: {text, size, bold, color, font}
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing

    for run_def in runs:
        r = p.add_run()
        r.text = run_def["text"]
        r.font.size = Pt(run_def.get("size", 18))
        r.font.bold = run_def.get("bold", False)
        r.font.color.rgb = run_def.get("color", ENCRE)
        r.font.name = run_def.get("font", FONT_BODY)
    return tb


def bullets(slide, x, y, w, h, items, *, size=14, color=INK, bullet_char="-",
            line_spacing=1.4, bold_first=False):
    """Multi-line bullet list."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)

        r = p.add_run()
        r.text = f"{bullet_char}  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = FONT_BODY
        if bold_first and i == 0:
            r.font.bold = True
    return tb


def numbered(slide, x, y, w, h, items, *, size=14, color=INK, line_spacing=1.5):
    """1. 2. 3. numbered list with colored numbers."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)

        # Number in accent color
        r_num = p.add_run()
        r_num.text = f"{i + 1}.  "
        r_num.font.size = Pt(size)
        r_num.font.bold = True
        r_num.font.color.rgb = ACCENT
        r_num.font.name = FONT_BODY

        # Body text
        r_body = p.add_run()
        r_body.text = item
        r_body.font.size = Pt(size)
        r_body.font.color.rgb = color
        r_body.font.name = FONT_BODY
    return tb


def rect(slide, x, y, w, h, fill=PAPIER_CHAUD, line=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if line_w is not None:
            sh.line.width = line_w
    return sh


def rounded(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.adjustments[0] = 0.25
    return sh


def hairline(slide, x, y, w, color=RULE_LIGHT):
    return rect(slide, x, y, w, Emu(9525), fill=color)  # 1 pixel-ish


def placeholder(slide, x, y, w, h, label):
    """Gray box telling the user a screenshot goes here."""
    rect(slide, x, y, w, h, fill=PAPIER_CHAUD, line=RULE_LIGHT, line_w=Pt(0.5))
    # Camera icon-ish text
    text(slide, x, y, w, h,
         f"⎘   CAPTURE À INSERER\n{label}",
         size=10, color=INK_DIM, font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         line_spacing=1.5)


def safe_image(slide, path, x, y, w, h, fallback_label="IMAGE"):
    full = os.path.join(ASSETS, path) if not os.path.isabs(path) else path
    if os.path.exists(full):
        slide.shapes.add_picture(full, x, y, w, h)
    else:
        placeholder(slide, x, y, w, h, fallback_label)


def footer(slide, page_num, total=18):
    """Mono footer with brand mark + pagination."""
    # left
    text(slide, Inches(0.7), Inches(7.05), Inches(7), Inches(0.3),
         "CARNET PLEIN® BY MAD MAKERS  ·  BONUS #1 · FICHE GOOGLE PARFAITE",
         size=8, color=INK_DIM, font=FONT_MONO, align=PP_ALIGN.LEFT)
    # right
    text(slide, Inches(11.6), Inches(7.05), Inches(1.1), Inches(0.3),
         f"{page_num:02d} / {total:02d}",
         size=8, color=INK_DIM, font=FONT_MONO, align=PP_ALIGN.RIGHT)


def accent_dot(slide, x, y, size=Inches(0.12)):
    """Square brand-mark accent."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()


def eyebrow(slide, x, y, w, label):
    """Mono uppercase eyebrow with orange square."""
    accent_dot(slide, x, y + Inches(0.05))
    text(slide, x + Inches(0.22), y, w, Inches(0.3),
         label, size=10, color=ACCENT, font=FONT_MONO,
         align=PP_ALIGN.LEFT)


# ===========================================================
# SLIDE 1 - COVER
# ===========================================================
def slide_cover():
    s = blank_slide(ENCRE)

    # Spotlight halo (soft orange glow upper-right)
    glow = slide_oval = s.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7.5), Inches(-2),
        Inches(8), Inches(8))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x2d, 0x18, 0x0c)  # dark orange tint
    glow.line.fill.background()

    # Badge top-left
    rounded(s, Inches(0.7), Inches(0.6), Inches(4.4), Inches(0.42), ACCENT)
    text(s, Inches(0.7), Inches(0.6), Inches(4.4), Inches(0.42),
         "BONUS #1  ·  CARNET PLEIN® BY MAD MAKERS",
         size=11, bold=True, color=ENCRE, font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Title
    text(s, Inches(0.7), Inches(2.2), Inches(11), Inches(2.4),
         "Fiche Google\nParfaite.",
         size=88, bold=True, color=WHITE, font=FONT_DISPLAY,
         line_spacing=0.95)

    # Subtitle accent line
    rect(s, Inches(0.7), Inches(5.0), Inches(0.6), Inches(0.04), fill=ACCENT)

    # Subtitle
    text(s, Inches(0.7), Inches(5.1), Inches(10), Inches(0.6),
         "Les 12 points qui font passer votre fiche en page 1 locale.",
         size=22, color=WHITE, font=FONT_BODY, line_spacing=1.25)

    # Tag chips
    chip_y = Inches(6.0)
    chips = [
        "12 actions concretes",
        "4 a 5h sur 30 jours",
        "Resultats visibles a partir de J21",
    ]
    cx = Inches(0.7)
    for ch in chips:
        # measure roughly by string length
        cw = Inches(0.18 * len(ch) + 0.6)
        rounded(s, cx, chip_y, cw, Inches(0.4),
                RGBColor(0x22, 0x22, 0x22))
        text(s, cx, chip_y, cw, Inches(0.4), ch.upper(),
             size=9, color=PAPER_DIM, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw + Inches(0.1)

    # Bottom-right brand mark
    text(s, Inches(8.5), Inches(7.05), Inches(4.2), Inches(0.3),
         "carnetplein.mad-makers.fr",
         size=10, color=PAPER_DIM, font=FONT_MONO,
         align=PP_ALIGN.RIGHT)


# ===========================================================
# SLIDE 2 - INTRO / POURQUOI
# ===========================================================
def slide_intro():
    s = blank_slide(PAPIER)
    eyebrow(s, Inches(0.7), Inches(0.55), Inches(6), "POURQUOI CETTE CHECKLIST")

    text(s, Inches(0.7), Inches(1.05), Inches(11.5), Inches(1.6),
         "76% des recherches locales sur mobile\ndebouchent sur un appel ou une visite.",
         size=40, bold=True, color=ENCRE, font=FONT_DISPLAY,
         line_spacing=1.05)

    text(s, Inches(0.7), Inches(2.95), Inches(7.5), Inches(2.5),
         "Quand un particulier tape « plombier + sa ville » sur Google, trois "
         "fiches s'affichent sur la carte avant tout le reste. C'est le bloc local "
         "(3-pack). Ces 3 fiches captent a elles seules environ 42% des clics organiques "
         "sur ce type de recherche (Backlinko, 2024). Si votre fiche n'apparait pas dans "
         "le 3-pack, vous laissez la quasi-totalite des prospects a vos concurrents.\n\n"
         "Appliquez les 12 points ci-dessous, votre visibilite monte et les appels suivent.",
         size=14, color=INK, font=FONT_BODY, line_spacing=1.55)

    # Right column: how to use
    rounded(s, Inches(8.6), Inches(2.95), Inches(4.05), Inches(3.6), PAPIER_CHAUD)
    eyebrow(s, Inches(8.85), Inches(3.15), Inches(4), "COMMENT UTILISER CE DOCUMENT")

    items = [
        "Travaillez les 12 points dans l'ordre - ils sont classes par impact",
        "Chaque point prend 10 a 45 minutes - etalez sur la semaine",
        "Cochez les cases au fur et a mesure. Recommencez tous les 6 mois.",
    ]
    bullets(s, Inches(8.85), Inches(3.7), Inches(3.65), Inches(2.7),
            items, size=12, color=INK, bullet_char="-")

    # Bottom callout 2026
    rounded(s, Inches(0.7), Inches(5.85), Inches(11.95), Inches(1.0),
            RGBColor(0x1a, 0x1c, 0x18))
    accent_dot(s, Inches(0.95), Inches(6.18))
    text(s, Inches(1.2), Inches(6.0), Inches(2.4), Inches(0.35),
         "NOTE 2026",
         size=10, bold=True, color=ACCENT, font=FONT_MONO)
    text(s, Inches(1.2), Inches(6.35), Inches(11.2), Inches(0.5),
         "Google ne gere plus les fiches mono-etablissement sur business.google.com. Tout passe par Google Search - tapez le nom "
         "de votre entreprise dans Google, le panneau de gestion s'affiche avec le bouton Modifier le profil.",
         size=11, color=WHITE, font=FONT_BODY, line_spacing=1.4)

    footer(s, 2)


# ===========================================================
# SLIDES 3-14 - LES 12 POINTS
# ===========================================================
POINTS = [
    {
        "num": "01",
        "title": "Choisir la bonne categorie principale",
        "why": ("Quand quelqu'un tape « plombier reims » ou « chauffagiste creil », "
                "Google regarde votre categorie principale en premier. Si elle est mauvaise "
                "ou floue, Google ne sait pas ou vous ranger. Vous etes invisible. "
                "C'est la base."),
        "actions": [
            "Tapez le nom de votre entreprise dans Google, connecte a votre compte proprietaire",
            "Votre fiche apparait, cliquez sur Modifier le profil",
            "A propos > Categorie principale",
            "Tapez votre metier exact : Plombier ou Plombier-chauffagiste. Pas Artisan, pas Entreprise du batiment",
            "Sauvegardez",
        ],
        "placeholder": "Editeur Categorie principale - GBP",
    },
    {
        "num": "02",
        "title": "Ajouter toutes vos categories secondaires",
        "why": ("Google vous laisse 1 categorie principale + jusqu'a 9 categories secondaires. "
                "Chaque categorie cochee vous rend visible sur une recherche associee. Si vous "
                "installez des PAC mais que Installateur de pompes a chaleur n'est pas cochee, "
                "vous ne ressortez pas sur cette recherche. Dommage, c'est la que se trouve le panier moyen le plus gros."),
        "actions": [
            "Modifier le profil > A propos > Categories supplementaires",
            "Cliquez sur Ajouter une autre categorie",
            "Cochez vos vrais services : Chauffagiste, Installateur de pompes a chaleur, Service de plomberie d'urgence...",
            "Uniquement ce que vous faites vraiment, sinon Google declasse votre fiche",
            "Sauvegardez",
        ],
        "placeholder": "Liste des categories secondaires - GBP",
    },
    {
        "num": "03",
        "title": "Ecrire une description avec vos vrais mots-cles",
        "why": ("Google lit votre description. Si vous y mettez « plombier-chauffagiste a "
                "Compiegne, specialiste pompe a chaleur et MaPrimeRenov », Google sait sur quoi "
                "vous faire ressortir. 750 caracteres maximum, seuls les 250 premiers sont "
                "visibles avant le Voir plus. Concentrez-y l'essentiel."),
        "actions": [
            "Modifier le profil > A propos > Description",
            "Premiers 250 caracteres : metier + ville principale + 1-2 specialites fortes",
            "Suite : zone d'intervention + certifications + offre devis",
            "Exemple : Plombier-chauffagiste a Reims, specialiste pompes a chaleur. Entreprise RGE QualiPAC...",
            "Pas de superlatifs (le meilleur, n°1). Restez factuel",
        ],
        "placeholder": "Editeur de Description - GBP",
    },
    {
        "num": "04",
        "title": "Charger un maximum de photos, rangees par type",
        "why": ("Les photos generent 30 a 50% de vues en plus et nourrissent les signaux d'engagement "
                "(clics, appels, demandes d'itineraire) que Google utilise pour mesurer la popularite. "
                "La moyenne secteur plomberie en France : 70 photos par fiche. Top 3 : 100+. "
                "Visez 30 minimum pour demarrer, 70 pour la cible."),
        "actions": [
            "Modifier le profil > Photos",
            "1 logo carre + 1 photo de couverture (camion ou equipe)",
            "5-10 photos de l'equipe, 3-5 du vehicule floque, 15-25 chantiers terminés",
            "10-15 installations propres : PAC, chaudiere, sanitaire",
            "Ajoutez 2-5 photos chaque semaine. La fraicheur compte",
        ],
        "placeholder": "Onglet Photos - GBP",
        "image": "metiers/electrician-fixing.png",
    },
    {
        "num": "05",
        "title": "Verrouiller les horaires, y compris jours feries",
        "why": ("Une fiche aux horaires obsoletes sort automatiquement des recherches filtrees "
                "« Ouvert maintenant ». Pire : un client qui se deplace pour vous trouver ferme "
                "un 1er mai laisse un avis 1 etoile en rentrant. Et un avis 1 etoile efface "
                "10 avis 5 etoiles chez le prospect suivant."),
        "actions": [
            "Modifier le profil > Horaires",
            "Renseignez vos horaires habituels jour par jour",
            "Horaires speciaux : ajoutez chaque jour ferie (1er janvier, lundi de Paques, 1er mai, 8 mai, Ascension, etc.)",
            "Depannage 24/7 : activez Ouvert 24h/24",
            "Repassez dessus en janvier et en juin chaque annee",
        ],
        "placeholder": "Editeur Horaires + Horaires speciaux - GBP",
    },
    {
        "num": "06",
        "title": "Mettre un vrai numero de telephone direct",
        "why": ("Le numero affiche doit decrocher. Pas un standard qui filtre, pas un portable "
                "qui ne repond pas en chantier. Google compte les appels recus depuis votre fiche : "
                "c'est un signal d'activite. Et un appel manque, c'est un prospect perdu, point."),
        "actions": [
            "Modifier le profil > Coordonnees > Telephone",
            "Mettez le numero qui decroche vraiment en journee",
            "Activez Numero supplementaire pour un fixe d'agence si besoin",
            "Numero identique sur site, cartes de visite, PagesJaunes, Habitatpresto",
            "La coherence Nom-Adresse-Telephone (NAP) reste un signal local confirme",
        ],
        "placeholder": "Champ Telephone - GBP",
    },
    {
        "num": "07",
        "title": "Delimiter clairement votre zone d'intervention",
        "why": ("Google a besoin de savoir ou vous travaillez pour vous proposer sur les bonnes "
                "recherches geographiques. Sans zone definie, vous ressortez uniquement sur votre "
                "commune. En la definissant, vous couvrez tous les villages alentours."),
        "actions": [
            "Modifier le profil > Etablissement et zone de chalandise",
            "Zones desservies > Ajouter une zone",
            "Saisissez vos departements : Marne, Aisne, Ardennes (par exemple)",
            "Ajoutez les 10-15 villes principales que vous couvrez",
            "N'inventez pas. Si vous ne vous deplacez pas a Sedan, ne la cochez pas",
        ],
        "placeholder": "Editeur Zones desservies - GBP",
    },
    {
        "num": "08",
        "title": "Lister chaque service avec sa propre description",
        "why": ("Chaque service ajoute est un mot-cle supplementaire pour Google, confirme comme "
                "facteur de ranking (impact modere, WebFX 2026). Une fiche avec 20 services bien "
                "decrits ressort sur 20 recherches differentes. Sans services detailles, vous "
                "ressortez sur 3."),
        "actions": [
            "Modifier le profil > Services",
            "Ajouter un autre service pour chaque prestation : Installation PAC air-eau, Depannage chaudiere, Renovation salle de bain...",
            "Pour chaque service : 2-3 lignes de description avec ville et certification (RGE, QualiPAC, Qualibat)",
            "Ajoutez un prix indicatif et un temps d'intervention estime quand c'est possible",
            "Sauvegardez chaque service",
        ],
        "placeholder": "Liste des Services - GBP",
    },
    {
        "num": "09",
        "title": "Publier une mise a jour chaque semaine",
        "why": ("Publier n'a pas d'effet direct sur le ranking mais une fiche qui publie recoit "
                "plus de vues et plus d'interactions - ces signaux d'engagement comptent pour 9 a "
                "11% du classement local (Whitespark 2026). C'est rapide a faire depuis votre "
                "telephone, entre deux chantiers."),
        "actions": [
            "Modifier le profil > Ajouter une mise a jour (icone + sur mobile)",
            "Hub Publications : Nouveautes / Offres / Evenements depuis 2025",
            "1 theme par semaine : chantier termine, conseil saisonnier, aide financiere (MaPrimeRenov)",
            "1 photo + 100-200 mots suffisent",
            "Bloquez un creneau fixe le vendredi matin",
        ],
        "placeholder": "Hub Publications - GBP",
    },
    {
        "num": "10",
        "title": "Demander un avis apres chaque chantier",
        "why": ("Les avis pesent 16 a 20% du classement Local Pack (Whitespark 2026) : c'est le "
                "2e facteur apres votre fiche elle-meme. Un client content ne laisse pas d'avis "
                "spontanement, il faut le demander. C'est chiant a mettre en place, c'est ce qui "
                "debloque le plus de prospects sur la duree."),
        "actions": [
            "Modifier le profil > Avis > Demandez plus d'avis (ordinateur uniquement, pas mobile)",
            "Google genere un lien court personnalise (g.page/r/... ou g.co/r/...) + QR code",
            "Demandez le jour meme de la fin de chantier, pas 3 semaines apres",
            "SMS type : Bonjour M. Dupont, content que l'installation soit OK. Si vous avez 30 secondes, un avis Google nous aiderait : [lien]",
            "Imprimez le QR code sur vos factures",
        ],
        "placeholder": "Bouton Demandez plus d'avis - GBP",
    },
    {
        "num": "11",
        "title": "Repondre a tous les avis",
        "why": ("88% des consommateurs preferent les entreprises qui repondent a tous leurs avis "
                "(BrightLocal 2024). Les fiches qui repondent a plus de 80% beneficient d'un boost "
                "mesurable. Surtout sur les avis negatifs : un avis 1 etoile sans reponse fait fuir, "
                "le meme avec une reponse calme et factuelle rassure."),
        "actions": [
            "Modifier le profil > Avis",
            "Delais : 48-72h pour un avis negatif, 1 semaine pour un positif",
            "Reponse positif en 2 lignes : remerciez, mentionnez le type de prestation",
            "Avis negatif : Bonjour [prenom], merci pour votre retour. Nous prenons note de votre deception. Pouvez-vous nous contacter au [numero] pour qu'on regarde ensemble ?",
            "Ne supprimez pas, ne contredisez pas publiquement. 15 min chaque lundi matin",
        ],
        "placeholder": "Onglet Reponses aux avis - GBP",
    },
    {
        "num": "12",
        "title": "Questions et reponses : alimentez votre site, pas la fiche",
        "why": ("Mise a jour fin 2025 : Google retire progressivement la section Questions/Reponses "
                "de la fiche. L'API a ete coupee le 3 novembre 2025, remplacee par une reponse IA "
                "Gemini (Ask Maps) qui pioche dans votre fiche, vos avis et votre site. "
                "Conclusion : nourrissez les sources que l'IA va lire."),
        "actions": [
            "Si Q&R encore visibles sur votre fiche : repondez aux questions deja posees, supprimez les reponses erronees",
            "Sur votre site : creez une page FAQ avec balisage Schema.org/FAQPage",
            "8-10 questions classiques : Etes-vous certifie RGE ? Faites-vous des devis gratuits ? Acceptez-vous MaPrimeRenov ?",
            "Remontez les questions frequentes dans la Description et les Services de la fiche",
            "C'est cette information structuree que Gemini lira pour repondre a vos prospects",
        ],
        "placeholder": "Section Q&R en retrait + FAQ Schema.org",
    },
]


def slide_point(point, page_num):
    s = blank_slide(PAPIER)

    # Big number left side
    text(s, Inches(0.7), Inches(0.55), Inches(1.5), Inches(0.4),
         f"{point['num']} / 12",
         size=12, bold=True, color=ACCENT, font=FONT_MONO)

    # Title
    text(s, Inches(0.7), Inches(1.0), Inches(11.5), Inches(1.0),
         point["title"],
         size=32, bold=True, color=ENCRE, font=FONT_DISPLAY, line_spacing=1.1)

    # Decorative orange line
    rect(s, Inches(0.7), Inches(2.1), Inches(0.5), Inches(0.04), fill=ACCENT)

    # 2-column body: Pourquoi (left) + Action (right)
    col1_x = Inches(0.7)
    col2_x = Inches(6.85)
    col_w = Inches(5.8)
    col_y = Inches(2.5)

    # Left col header
    text(s, col1_x, col_y, col_w, Inches(0.3),
         "POURQUOI C'EST IMPORTANT",
         size=10, bold=True, color=ACCENT, font=FONT_MONO)

    # Left col body
    text(s, col1_x, col_y + Inches(0.45), col_w, Inches(3),
         point["why"],
         size=13, color=INK, font=FONT_BODY, line_spacing=1.55)

    # Right col header
    text(s, col2_x, col_y, col_w, Inches(0.3),
         "ACTION CONCRETE",
         size=10, bold=True, color=ACCENT, font=FONT_MONO)

    # Right col actions (numbered)
    numbered(s, col2_x, col_y + Inches(0.45), col_w, Inches(3.5),
             point["actions"], size=12, color=INK, line_spacing=1.5)

    # Checkbox + placeholder area at bottom
    checkbox_y = Inches(6.4)

    # Checkbox
    rect(s, Inches(0.7), checkbox_y, Inches(0.28), Inches(0.28),
         fill=PAPIER, line=INK, line_w=Pt(1.2))
    text(s, Inches(1.05), checkbox_y - Inches(0.04), Inches(2), Inches(0.4),
         "Fait",
         size=12, color=INK, font=FONT_BODY)

    # Placeholder / Image on the right at the bottom
    ph_w = Inches(5.0)
    ph_h = Inches(0.55)
    ph_x = Inches(7.65)
    ph_y = checkbox_y - Inches(0.05)

    if "image" in point:
        safe_image(s, point["image"], ph_x, ph_y, ph_w, ph_h, point["placeholder"])
    else:
        placeholder(s, ph_x, ph_y, ph_w, ph_h, point["placeholder"])

    footer(s, page_num)


# ===========================================================
# SLIDE 15 - PLAN D'ACTION
# ===========================================================
def slide_plan():
    s = blank_slide(PAPIER)
    eyebrow(s, Inches(0.7), Inches(0.55), Inches(6), "PLAN D'ACTION")
    text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(1.0),
         "Si vous n'avez le temps que pour 3 points cette semaine.",
         size=32, bold=True, color=ENCRE, font=FONT_DISPLAY, line_spacing=1.05)

    # Top 3 priority callout
    rounded(s, Inches(0.7), Inches(2.2), Inches(11.95), Inches(1.1), ACCENT)
    text(s, Inches(1.0), Inches(2.35), Inches(11.5), Inches(0.4),
         "PRIORITES",
         size=10, bold=True, color=ENCRE, font=FONT_MONO)
    text(s, Inches(1.0), Inches(2.7), Inches(11.5), Inches(0.5),
         "Points 01, 04 et 10. Bonne categorie + 70 photos visees + systeme de demande d'avis post-chantier = environ 80% du resultat. "
         "Le reste affine et solidifie.",
         size=13, color=ENCRE, font=FONT_BODY, line_spacing=1.4)

    # 30-day plan title
    text(s, Inches(0.7), Inches(3.7), Inches(11), Inches(0.5),
         "Plan sur 30 jours.",
         size=22, bold=True, color=ENCRE, font=FONT_DISPLAY)

    # 4 week boxes
    weeks = [
        ("S1 - FONDATIONS", "1h30",
         "Lundi  -  point 01\nMercredi  -  point 02\nVendredi  -  point 03"),
        ("S2 - VISIBILITE", "2h",
         "Lundi  -  point 04\nMercredi  -  point 05\nVendredi  -  point 06"),
        ("S3 - PERIMETRE", "2h",
         "Lundi  -  point 07\nMercredi  -  point 08\nVendredi  -  point 09"),
        ("S4 - TRACTION", "1h30",
         "Lundi  -  point 10\nMercredi  -  point 11\nVendredi  -  point 12"),
    ]

    box_w = Inches(2.85)
    box_h = Inches(2.15)
    gap = Inches(0.15)
    start_x = Inches(0.7)
    start_y = Inches(4.4)

    for i, (label, dur, content) in enumerate(weeks):
        bx = start_x + (box_w + gap) * i
        rounded(s, bx, start_y, box_w, box_h, PAPIER_CHAUD)
        text(s, bx + Inches(0.25), start_y + Inches(0.2),
             box_w - Inches(0.5), Inches(0.3),
             label, size=10, bold=True, color=ACCENT, font=FONT_MONO)
        text(s, bx + Inches(0.25), start_y + Inches(0.55),
             box_w - Inches(0.5), Inches(0.3),
             dur + " environ", size=10, color=INK_DIM, font=FONT_MONO)
        text(s, bx + Inches(0.25), start_y + Inches(1.0),
             box_w - Inches(0.5), Inches(1.0),
             content, size=12, color=INK, font=FONT_BODY, line_spacing=1.55)

    footer(s, 15)


# ===========================================================
# SLIDE 16 - ALERTES 2026
# ===========================================================
def slide_2026():
    s = blank_slide(ENCRE)
    eyebrow(s, Inches(0.7), Inches(0.55), Inches(6), "A GARDER EN TETE EN 2026")
    text(s, Inches(0.7), Inches(1.05), Inches(11.5), Inches(1.0),
         "Deux changements reglementaires qui changent la donne.",
         size=30, bold=True, color=WHITE, font=FONT_DISPLAY, line_spacing=1.1)

    # Two cards side by side
    card_w = Inches(5.9)
    card_h = Inches(3.8)
    card_y = Inches(2.5)

    # Left: RGE renforce
    rounded(s, Inches(0.7), card_y, card_w, card_h,
            RGBColor(0x1a, 0x1c, 0x18))
    text(s, Inches(0.95), card_y + Inches(0.3), Inches(5.5), Inches(0.3),
         "LABEL RGE RENFORCE",
         size=10, bold=True, color=ACCENT, font=FONT_MONO)
    text(s, Inches(0.95), card_y + Inches(0.65), Inches(5.5), Inches(0.5),
         "Depuis la loi du 30 juin 2025",
         size=12, color=PAPER_DIM, font=FONT_BODY)
    text(s, Inches(0.95), card_y + Inches(1.25), Inches(5.5), Inches(2.4),
         "Nouvelle obligation d'information du client, motif de retrait du label etendu. "
         "Argument cle dans le combat contre la fraude aux aides publiques.\n\n"
         "Mentionnez vos certifications RGE QualiPAC et Qualibat sur la fiche, dans la "
         "description, dans les services, et dans les posts hebdomadaires.",
         size=12, color=WHITE, font=FONT_BODY, line_spacing=1.6)

    # Right: MaPrimeRenov 2025
    rounded(s, Inches(6.75), card_y, card_w, card_h,
            RGBColor(0x1a, 0x1c, 0x18))
    text(s, Inches(7.0), card_y + Inches(0.3), Inches(5.5), Inches(0.3),
         "MAPRIMERENOV' 2025",
         size=10, bold=True, color=ACCENT, font=FONT_MONO)
    text(s, Inches(7.0), card_y + Inches(0.65), Inches(5.5), Inches(0.5),
         "Decret du 4 decembre 2024",
         size=12, color=PAPER_DIM, font=FONT_BODY)
    text(s, Inches(7.0), card_y + Inches(1.25), Inches(5.5), Inches(2.4),
         "Avance de tresorerie reduite : jusqu'a 50% de la prime versee avant chantier, "
         "reservee aux menages tres modestes (categorie bleue) en mono-geste, et aux "
         "bleus + jaunes en renovation d'ampleur.\n\n"
         "Argument fort a mettre en avant dans vos posts et votre FAQ site.",
         size=12, color=WHITE, font=FONT_BODY, line_spacing=1.6)

    footer(s, 16)


# ===========================================================
# SLIDE 17 - CARNET PLEIN(R) CTA
# ===========================================================
def slide_cta():
    s = blank_slide(PAPIER)

    # Top accent line
    rect(s, Inches(0.7), Inches(0.7), Inches(0.5), Inches(0.04), fill=ACCENT)
    text(s, Inches(0.7), Inches(0.85), Inches(11), Inches(0.4),
         "ALLER PLUS LOIN", size=10, bold=True, color=ACCENT, font=FONT_MONO)

    # Headline
    text(s, Inches(0.7), Inches(1.5), Inches(11.5), Inches(2.2),
         "Cette checklist, on l'applique a chaque artisan signe du programme.",
         size=38, bold=True, color=ENCRE, font=FONT_DISPLAY, line_spacing=1.1)

    # Body
    text(s, Inches(0.7), Inches(4.0), Inches(7.5), Inches(2.5),
         "Plus l'inscription sur 50 annuaires metier (PagesJaunes, Habitatpresto, Travaux.com...), "
         "un systeme automatise qui envoie la demande d'avis au bon moment apres chaque "
         "chantier, et un reporting mensuel clair.\n\n"
         "Vous gardez le savoir-faire metier. On gere le digital.",
         size=15, color=INK, font=FONT_BODY, line_spacing=1.6)

    # Right CTA card
    cta_x = Inches(8.6)
    cta_y = Inches(4.0)
    cta_w = Inches(4.05)
    cta_h = Inches(2.5)
    rounded(s, cta_x, cta_y, cta_w, cta_h, ENCRE)
    text(s, cta_x + Inches(0.3), cta_y + Inches(0.3),
         cta_w - Inches(0.6), Inches(0.3),
         "AUDIT GRATUIT 20 MIN",
         size=10, bold=True, color=ACCENT, font=FONT_MONO)
    text(s, cta_x + Inches(0.3), cta_y + Inches(0.75),
         cta_w - Inches(0.6), Inches(1.4),
         "On regarde votre marche local et on vous dit honnetement si Carnet Plein® est fait pour vous. Sans engagement, sans pitch commercial.",
         size=12, color=WHITE, font=FONT_BODY, line_spacing=1.5)

    # CTA button
    rounded(s, cta_x + Inches(0.3), cta_y + Inches(1.85),
            cta_w - Inches(0.6), Inches(0.5), ACCENT)
    text(s, cta_x + Inches(0.3), cta_y + Inches(1.85),
         cta_w - Inches(0.6), Inches(0.5),
         "PRO.MAD-MAKERS.FR",
         size=11, bold=True, color=WHITE, font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 17)


# ===========================================================
# SLIDE 18 - SOURCES
# ===========================================================
def slide_sources():
    s = blank_slide(PAPIER)
    eyebrow(s, Inches(0.7), Inches(0.55), Inches(6), "SOURCES PRINCIPALES")
    text(s, Inches(0.7), Inches(1.05), Inches(11), Inches(1.0),
         "Les references derriere les chiffres.",
         size=28, bold=True, color=ENCRE, font=FONT_DISPLAY, line_spacing=1.1)

    sources = [
        "Aide Google Business Profile  ·  support.google.com/business",
        "Whitespark Local Search Ranking Factors 2026",
        "BrightLocal Local Consumer Review Survey 2024-2026",
        "Backlinko Local SEO Statistics 2024",
        "Sterling Sky (Google Posts ranking factor study)",
        "ANAH / monprojet.anah.gouv.fr (MaPrimeRenov')",
        "Qualit'EnR (label RGE)",
        "economie.gouv.fr - Loi 2025-594 du 30 juin 2025",
        "WebFX Local SEO Statistics 2026",
    ]

    bullets(s, Inches(0.7), Inches(2.5), Inches(11), Inches(4),
            sources, size=14, color=INK, bullet_char="-",
            line_spacing=1.7)

    # Bottom legal-ish note
    text(s, Inches(0.7), Inches(6.5), Inches(11.5), Inches(0.5),
         "Document a usage interne du client. Donnees publiques chiffrees actualisees au 1er trimestre 2026. "
         "Toute reproduction ou diffusion sans autorisation est interdite.",
         size=9, color=INK_DIM, font=FONT_BODY, line_spacing=1.4)

    footer(s, 18)


# ===========================================================
# BUILD
# ===========================================================
slide_cover()
slide_intro()

for idx, point in enumerate(POINTS):
    slide_point(point, page_num=3 + idx)

slide_plan()
slide_2026()
slide_cta()
slide_sources()

out_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "01-fiche-google-parfaite.pptx"
)
prs.save(out_path)
print(f"OK - {len(prs.slides)} slides written to:")
print(f"     {out_path}")
print(f"     size: {os.path.getsize(out_path):,} bytes")
